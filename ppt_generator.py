"""
PPT生成引擎
---------
用于将分析后的内容转化为结构化的PPT演示文稿,
包括内容组织、模板选择、图形设计等功能。
"""

import os
import json
import random
import re
from typing import List, Dict, Tuple, Optional, Union, Any
import numpy as np
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt
import networkx as nx
from io import BytesIO
import base64
import logging

# 尝试导入pptx库
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx库未安装，PPT生成功能将受限")

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class PPTGenerator:
    """PPT生成引擎主类"""
    
    def __init__(self, config: Dict[str, Any] = None):
        """
        初始化PPT生成引擎
        
        Args:
            config: 配置参数，包括模板设置、样式选项等
        """
        self.config = config or {}
        
        # 检查必要库
        if not PPTX_AVAILABLE:
            raise ImportError("PPT生成需要python-pptx库，请使用pip install python-pptx安装")
        
        # 加载模板
        self.template_path = self.config.get('template_path')
        self.template = self._load_template()
        
        # 设置默认样式
        self.styles = self._load_styles()
        
        # 默认图表颜色
        self.chart_colors = self.config.get('chart_colors', [
            '#4472C4', '#ED7D31', '#A5A5A5', '#FFC000', 
            '#5B9BD5', '#70AD47', '#264478', '#9E480E'
        ])
    
    def _load_template(self) -> Presentation:
        """加载PPT模板"""
        if self.template_path and os.path.exists(self.template_path):
            return Presentation(self.template_path)
        else:
            # 使用默认空白模板
            return Presentation()
    
    def _load_styles(self) -> Dict[str, Any]:
        """加载样式配置"""
        # 从配置文件加载或使用默认样式
        default_styles = {
            'title_font': {
                'name': '微软雅黑',
                'size': Pt(40),
                'bold': True,
                'color': RGBColor(33, 33, 33)
            },
            'heading_font': {
                'name': '微软雅黑',
                'size': Pt(32),
                'bold': True,
                'color': RGBColor(33, 33, 33)
            },
            'subheading_font': {
                'name': '微软雅黑',
                'size': Pt(24),
                'bold': True,
                'color': RGBColor(33, 33, 33)
            },
            'body_font': {
                'name': '微软雅黑',
                'size': Pt(18),
                'bold': False,
                'color': RGBColor(33, 33, 33)
            },
            'theme_colors': {
                'primary': RGBColor(0, 112, 192),
                'secondary': RGBColor(0, 176, 80),
                'accent1': RGBColor(255, 192, 0),
                'accent2': RGBColor(192, 0, 0),
                'background': RGBColor(255, 255, 255)
            }
        }
        
        # 合并用户配置和默认配置
        styles = self.config.get('styles', {})
        for key, default_value in default_styles.items():
            if key not in styles:
                styles[key] = default_value
        
        return styles
    
    def generate_ppt(self, content_analysis: Dict[str, Any], output_path: str) -> str:
        """
        生成PPT演示文稿
        
        Args:
            content_analysis: 内容分析结果
            output_path: 输出的PPT文件路径
            
        Returns:
            生成的PPT文件路径
        """
        logger.info("开始生成PPT...")
        
        # 创建演示文稿
        prs = self.template
        
        # 提取学科和难度信息
        subject = content_analysis.get('subject', 'general')
        difficulty = content_analysis.get('difficulty_level', 'intermediate')
        
        # 添加封面
        self._add_cover_slide(prs, content_analysis)
        
        # 添加目录
        self._add_toc_slide(prs, content_analysis)
        
        # 添加章节内容
        chapters = content_analysis.get('chapters', [])
        for chapter in chapters:
            self._add_chapter_slides(prs, chapter, subject, difficulty)
        
        # 添加知识图谱
        if 'knowledge_graph' in content_analysis:
            self._add_knowledge_graph_slide(prs, content_analysis['knowledge_graph'])
        
        # 添加总结和测验
        if len(chapters) > 0:
            self._add_summary_slide(prs, content_analysis)
            self._add_quiz_slide(prs, content_analysis)
        
        # 保存PPT
        prs.save(output_path)
        logger.info(f"PPT已保存至: {output_path}")
        
        return output_path
    
    def _add_cover_slide(self, prs: Presentation, content_analysis: Dict[str, Any]) -> None:
        """添加封面幻灯片"""
        # 选择封面布局
        slide_layout = prs.slide_layouts[0]  # 通常第一个布局是标题幻灯片
        slide = prs.slides.add_slide(slide_layout)
        
        # 获取标题和副标题占位符
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        # 设置标题
        main_topics = content_analysis.get('main_topics', [])
        title_text = "教材内容讲解"
        if main_topics and len(main_topics) > 0:
            title_text = main_topics[0]  # 使用第一个主题作为标题
            
        title.text = title_text
        
        # 设置副标题
        if subtitle:
            subject = content_analysis.get('subject', '')
            if subject == 'physics':
                subject_text = "物理"
            elif subject == 'chemistry':
                subject_text = "化学"
            elif subject == 'biology':
                subject_text = "生物"
            elif subject == 'math':
                subject_text = "数学"
            else:
                subject_text = subject.capitalize()
                
            subtitle.text = f"{subject_text} - 交互式学习"
        
        # 应用样式
        self._apply_text_style(title.text_frame.paragraphs[0], 'title_font')
        if subtitle:
            self._apply_text_style(subtitle.text_frame.paragraphs[0], 'subheading_font')
    
    def _add_toc_slide(self, prs: Presentation, content_analysis: Dict[str, Any]) -> None:
        """添加目录幻灯片"""
        # 选择布局
        slide_layout = prs.slide_layouts[1]  # 通常第二个布局是标题和内容
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = "课程内容"
        self._apply_text_style(title.text_frame.paragraphs[0], 'heading_font')
        
        # 获取内容占位符
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # 内容占位符类型
                content_placeholder = shape
                break
        
        # 如果找不到内容占位符，创建文本框
        if not content_placeholder:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        # 添加目录内容
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # 清除默认文本
        
        chapters = content_analysis.get('chapters', [])
        for i, chapter in enumerate(chapters):
            title = chapter.get('title', f'章节 {i+1}')
            p = text_frame.add_paragraph()
            p.text = f"{i+1}. {title}"
            self._apply_text_style(p, 'body_font')
            p.space_after = Pt(12)
    
    def _add_chapter_slides(self, prs: Presentation, chapter: Dict[str, Any], 
                           subject: str, difficulty: str) -> None:
        """为每个章节添加幻灯片"""
        # 章节标题幻灯片
        self._add_chapter_title_slide(prs, chapter)
        
        # 章节内容幻灯片
        self._add_chapter_content_slides(prs, chapter, subject, difficulty)
        
        # 如果有学习目标，添加学习目标幻灯片
        if 'learning_objectives' in chapter and chapter['learning_objectives']:
            self._add_learning_objectives_slide(prs, chapter)
        
        # 如果有子章节，递归添加
        if 'subchapters' in chapter:
            for subchapter in chapter['subchapters']:
                self._add_subchapter_slides(prs, subchapter, subject, difficulty)
    
    def _add_chapter_title_slide(self, prs: Presentation, chapter: Dict[str, Any]) -> None:
        """添加章节标题幻灯片"""
        # 使用章节标题布局
        slide_layout = prs.slide_layouts[2]  # 通常是标题和内容，可以根据实际模板调整
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = chapter.get('title', '未命名章节')
        self._apply_text_style(title.text_frame.paragraphs[0], 'heading_font')
        
        # 如果有摘要，添加摘要
        if 'summary' in chapter and chapter['summary']:
            # 找到内容占位符
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # 内容占位符类型
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                content_placeholder.text = chapter['summary']
                self._apply_text_style(content_placeholder.text_frame.paragraphs[0], 'body_font')
    
    def _add_chapter_content_slides(self, prs: Presentation, chapter: Dict[str, Any], 
                                   subject: str, difficulty: str) -> None:
        """添加章节内容幻灯片"""
        # 获取章节关键概念
        key_concepts = chapter.get('key_concepts', [])
        
        if key_concepts:
            # 添加关键概念幻灯片
            self._add_key_concepts_slide(prs, chapter['title'], key_concepts)
        
        # 根据学科和难度选择合适的可视化方式
        if subject in ['physics', 'chemistry'] and key_concepts:
            # 为物理和化学添加概念关系图
            self._add_concept_relationship_slide(prs, chapter['title'], key_concepts)
        
        # 如果有足够的概念，创建知识点详解幻灯片
        if len(key_concepts) >= 3:
            # 每张幻灯片显示3个概念
            for i in range(0, len(key_concepts), 3):
                concepts_batch = key_concepts[i:i+3]
                self._add_concept_detail_slide(prs, chapter['title'], concepts_batch)
    
    def _add_key_concepts_slide(self, prs: Presentation, chapter_title: str, 
                               key_concepts: List[str]) -> None:
        """添加关键概念幻灯片"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = f"{chapter_title} - 关键概念"
        self._apply_text_style(title.text_frame.paragraphs[0], 'subheading_font')
        
        # 获取内容占位符
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # 内容占位符类型
                content_placeholder = shape
                break
        
        if not content_placeholder:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        # 添加关键概念列表
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for concept in key_concepts[:10]:  # 限制显示数量
            p = text_frame.add_paragraph()
            p.text = f"• {concept}"
            self._apply_text_style(p, 'body_font')
            p.space_after = Pt(12)
    
    def _add_concept_relationship_slide(self, prs: Presentation, chapter_title: str, 
                                       key_concepts: List[str]) -> None:
        """添加概念关系图幻灯片"""
        slide_layout = prs.slide_layouts[5]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_shape.text = f"{chapter_title} - 概念关系图"
        self._apply_text_style(title_shape.text_frame.paragraphs[0], 'subheading_font')
        
        # 生成概念关系图并插入
        img_path = self._generate_concept_graph(key_concepts, chapter_title)
        if img_path:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            slide.shapes.add_picture(img_path, left, top, width, height)
            
            # 删除临时文件
            try:
                os.remove(img_path)
            except:
                pass
    
    def _generate_concept_graph(self, concepts: List[str], title: str) -> Optional[str]:
        """生成概念关系图并返回图像文件路径"""
        if len(concepts) < 2:
            return None
            
        try:
            # 创建图
            G = nx.Graph()
            
            # 添加中心节点
            G.add_node(title, type='center')
            
            # 添加概念节点并连接到中心
            for concept in concepts:
                G.add_node(concept, type='concept')
                G.add_edge(title, concept)
            
            # 随机添加概念之间的关系
            num_extra_edges = min(len(concepts), 5)  # 限制额外边的数量
            for _ in range(num_extra_edges):
                c1, c2 = random.sample(concepts, 2)
                if not G.has_edge(c1, c2):
                    G.add_edge(c1, c2)
            
            # 设置节点颜色
            node_colors = []
            for node in G.nodes():
                if G.nodes[node].get('type') == 'center':
                    node_colors.append(self.chart_colors[0])
                else:
                    node_colors.append(self.chart_colors[1])
            
            # 创建图像
            plt.figure(figsize=(10, 8))
            pos = nx.spring_layout(G, seed=42)
            
            # 绘制节点和边
            nx.draw_networkx_nodes(G, pos, node_color=node_colors, node_size=2000, alpha=0.8)
            nx.draw_networkx_edges(G, pos, width=2, alpha=0.5)
            
            # 绘制标签
            nx.draw_networkx_labels(G, pos, font_size=12, font_family='SimHei')
            
            plt.axis('off')
            
            # 保存图像
            temp_img_path = 'temp_concept_graph.png'
            plt.savefig(temp_img_path, dpi=300, bbox_inches='tight')
            plt.close()
            
            return temp_img_path
            
        except Exception as e:
            logger.error(f"生成概念图失败: {e}")
            return None
    
    def _add_concept_detail_slide(self, prs: Presentation, chapter_title: str, 
                                 concepts: List[str]) -> None:
        """添加概念详解幻灯片"""
        slide_layout = prs.slide_layouts[3]  # 有多个内容部分的布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = f"{chapter_title} - 知识点详解"
        self._apply_text_style(title.text_frame.paragraphs[0], 'subheading_font')
        
        # 创建每个概念的详解
        left_positions = [Inches(0.5), Inches(4.0), Inches(7.5)]
        width = Inches(3.0)
        
        for i, concept in enumerate(concepts):
            if i >= 3:  # 限制每张幻灯片最多3个概念
                break
                
            # 创建一个文本框
            top = Inches(1.5)
            height = Inches(5)
            left = left_positions[i]
            
            txt_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = txt_box.text_frame
            
            # 添加概念标题
            p = text_frame.add_paragraph()
            p.text = concept
            self._apply_text_style(p, 'subheading_font')
            p.space_after = Pt(12)
            
            # 添加概念说明（模拟内容）
            p = text_frame.add_paragraph()
            p.text = f"这是关于{concept}的详细说明..."
            self._apply_text_style(p, 'body_font')
    
    def _add_subchapter_slides(self, prs: Presentation, subchapter: Dict[str, Any], 
                              subject: str, difficulty: str) -> None:
        """添加子章节幻灯片"""
        # 子章节结构和主章节类似，但使用不同的样式
        
        # 子章节标题幻灯片
        slide_layout = prs.slide_layouts[2]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = subchapter.get('title', '未命名子章节')
        self._apply_text_style(title.text_frame.paragraphs[0], 'subheading_font')
        
        # 如果有摘要，添加摘要
        if 'summary' in subchapter and subchapter['summary']:
            # 找到内容占位符
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # 内容占位符类型
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                content_placeholder.text = subchapter['summary']
                self._apply_text_style(content_placeholder.text_frame.paragraphs[0], 'body_font')
        
        # 处理子章节内容（如果有关键概念）
        key_concepts = subchapter.get('key_concepts', [])
        if key_concepts:
            self._add_key_concepts_slide(prs, subchapter['title'], key_concepts)
    
    def _add_learning_objectives_slide(self, prs: Presentation, chapter: Dict[str, Any]) -> None:
        """添加学习目标幻灯片"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = f"{chapter.get('title', '未命名章节')} - 学习目标"
        self._apply_text_style(title.text_frame.paragraphs[0], 'subheading_font')
        
        # 获取内容占位符
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # 内容占位符类型
                content_placeholder = shape
                break
        
        if not content_placeholder:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        # 添加学习目标
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        objectives = chapter.get('learning_objectives', [])
        for i, objective in enumerate(objectives):
            p = text_frame.add_paragraph()
            p.text = f"{i+1}. {objective}"
            self._apply_text_style(p, 'body_font')
            p.space_after = Pt(12)
    
    def _add_knowledge_graph_slide(self, prs: Presentation, knowledge_graph: Dict[str, Any]) -> None:
        """添加知识图谱幻灯片"""
        slide_layout = prs.slide_layouts[5]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_shape.text = "知识图谱"
        self._apply_text_style(title_shape.text_frame.paragraphs[0], 'heading_font')
        
        # 生成知识图谱可视化
        img_path = self._generate_knowledge_graph_image(knowledge_graph)
        if img_path:
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            slide.shapes.add_picture(img_path, left, top, width, height)
            
            # 删除临时文件
            try:
                os.remove(img_path)
            except:
                pass
    
    def _generate_knowledge_graph_image(self, knowledge_graph: Dict[str, Any]) -> Optional[str]:
        """生成知识图谱可视化图像"""
        try:
            # 创建有向图
            G = nx.DiGraph()
            
            # 添加节点
            nodes = knowledge_graph.get('nodes', [])
            for node in nodes:
                node_id = node.get('id', '')
                node_type = node.get('type', 'unknown')
                
                if node_id:
                    G.add_node(node_id, type=node_type)
            
            # 添加边
            edges = knowledge_graph.get('edges', [])
            for edge in edges:
                source = edge.get('source', '')
                target = edge.get('target', '')
                edge_type = edge.get('type', 'related')
                
                if source and target and source in G and target in G:
                    G.add_edge(source, target, type=edge_type)
            
            if not G.nodes():
                logger.warning("知识图谱为空，无法生成可视化")
                return None
            
            # 设置节点颜色
            node_colors = []
            for node in G.nodes():
                node_type = G.nodes[node].get('type', '')
                if node_type == 'chapter':
                    node_colors.append(self.chart_colors[0])
                elif node_type == 'concept':
                    node_colors.append(self.chart_colors[1])
                else:
                    node_colors.append(self.chart_colors[2])
            
            # 设置边颜色
            edge_colors = []
            for u, v, data in G.edges(data=True):
                edge_type = data.get('type', '')
                if edge_type == 'contains':
                    edge_colors.append(self.chart_colors[3])
                elif edge_type == 'related':
                    edge_colors.append(self.chart_colors[4])
                elif edge_type == 'next':
                    edge_colors.append(self.chart_colors[5])
                else:
                    edge_colors.append(self.chart_colors[6])
            
            # 创建可视化
            plt.figure(figsize=(12, 10))
            
            # 限制显示的节点数量
            if len(G.nodes()) > 30:
                # 截取图中最重要的部分
                central_nodes = [n for n, d in G.nodes(data=True) if d.get('type') == 'chapter'][:5]
                sub_nodes = []
                for node in central_nodes:
                    sub_nodes.append(node)
                    # 添加与中心节点相连的节点
                    sub_nodes.extend(list(G.successors(node))[:5])
                
                # 创建子图
                subG = G.subgraph(sub_nodes)
                pos = nx.spring_layout(subG, seed=42)
                
                # 绘制节点
                nx.draw_networkx_nodes(subG, pos, node_color=node_colors[:len(subG.nodes())], 
                                     node_size=700, alpha=0.8)
                
                # 绘制边
                nx.draw_networkx_edges(subG, pos, edge_color=edge_colors[:len(subG.edges())], 
                                     width=1.5, alpha=0.6, arrows=True, arrowsize=15)
                
                # 绘制标签
                nx.draw_networkx_labels(subG, pos, font_size=8, font_family="SimHei")
            else:
                pos = nx.spring_layout(G, seed=42)
                
                # 绘制节点
                nx.draw_networkx_nodes(G, pos, node_color=node_colors, node_size=700, alpha=0.8)
                
                # 绘制边
                nx.draw_networkx_edges(G, pos, edge_color=edge_colors, width=1.5, 
                                     alpha=0.6, arrows=True, arrowsize=15)
                
                # 绘制标签
                nx.draw_networkx_labels(G, pos, font_size=8, font_family="SimHei")
            
            plt.title("教材内容知识图谱")
            plt.axis('off')
            plt.tight_layout()
            
            # 保存图像
            temp_img_path = 'temp_knowledge_graph.png'
            plt.savefig(temp_img_path, dpi=300, bbox_inches='tight')
            plt.close()
            
            return temp_img_path
            
        except Exception as e:
            logger.error(f"生成知识图谱可视化失败: {e}")
            return None
    
    def _add_summary_slide(self, prs: Presentation, content_analysis: Dict[str, Any]) -> None:
        """添加总结幻灯片"""
        slide_layout = prs.slide_layouts[2]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = "课程总结"
        self._apply_text_style(title.text_frame.paragraphs[0], 'heading_font')
        
        # 获取内容占位符
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # 内容占位符类型
                content_placeholder = shape
                break
        
        if not content_placeholder:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        # 添加总结内容
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # 提取重要概念
        important_concepts = content_analysis.get('important_concepts', [])
        
        # 添加学习要点
        p = text_frame.add_paragraph()
        p.text = "课程要点："
        self._apply_text_style(p, 'subheading_font')
        
        # 添加重要概念
        for i, concept in enumerate(important_concepts[:5]):
            concept_name = concept.get('concept', '')
            if concept_name:
                p = text_frame.add_paragraph()
                p.text = f"• {concept_name}"
                self._apply_text_style(p, 'body_font')
                p.space_after = Pt(6)
        
        # 添加学习建议（如果有）
        teaching_suggestions = content_analysis.get('teaching_suggestions', [])
        if teaching_suggestions:
            p = text_frame.add_paragraph()
            p.text = "\n学习建议："
            self._apply_text_style(p, 'subheading_font')
            
            for i, suggestion in enumerate(teaching_suggestions[:3]):
                p = text_frame.add_paragraph()
                p.text = f"• {suggestion}"
                self._apply_text_style(p, 'body_font')
                p.space_after = Pt(6)
    
    def _add_quiz_slide(self, prs: Presentation, content_analysis: Dict[str, Any]) -> None:
        """添加测验幻灯片"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = "知识检测"
        self._apply_text_style(title.text_frame.paragraphs[0], 'heading_font')
        
        # 获取内容占位符
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # 内容占位符类型
                content_placeholder = shape
                break
        
        if not content_placeholder:
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)
            content_placeholder = slide.shapes.add_textbox(left, top, width, height)
        
        # 添加测验题目（这里只是示例，实际项目中应该根据内容动态生成）
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # 从重要概念生成测验题
        important_concepts = content_analysis.get('important_concepts', [])
        
        if important_concepts:
            quiz_concepts = important_concepts[:3]  # 取前3个重要概念
            
            for i, concept in enumerate(quiz_concepts):
                concept_name = concept.get('concept', '')
                if concept_name:
                    p = text_frame.add_paragraph()
                    p.text = f"问题 {i+1}: 请解释{concept_name}的含义及其重要性。"
                    self._apply_text_style(p, 'subheading_font')
                    p.space_after = Pt(12)
        else:
            # 如果没有重要概念，添加一个通用问题
            p = text_frame.add_paragraph()
            p.text = "问题: 请总结本课程中的主要知识点。"
            self._apply_text_style(p, 'subheading_font')
    
    def _apply_text_style(self, paragraph, style_name: str) -> None:
        """应用文本样式到段落"""
        style = self.styles.get(style_name)
        if not style:
            return
            
        # 应用字体样式
        for run in paragraph.runs:
            if 'name' in style:
                run.font.name = style['name']
            if 'size' in style:
                run.font.size = style['size']
            if 'bold' in style:
                run.font.bold = style['bold']
            if 'italic' in style:
                run.font.italic = style['italic']
            if 'color' in style:
                run.font.color.rgb = style['color']

def main():
    """主函数"""
    import argparse
    
    parser = argparse.ArgumentParser(description='生成教材内容PPT')
    parser.add_argument('input_file', help='输入的分析结果JSON文件')
    parser.add_argument('--output', '-o', help='输出的PPT文件路径')
    parser.add_argument('--template', '-t', help='PPT模板文件路径')
    
    args = parser.parse_args()
    
    # 加载内容分析结果
    with open(args.input_file, 'r', encoding='utf-8') as f:
        content_analysis = json.load(f)
    
    # 配置
    config = {}
    if args.template:
        config['template_path'] = args.template
    
    # 设置默认输出路径
    output_path = args.output
    if not output_path:
        input_base = os.path.splitext(os.path.basename(args.input_file))[0]
        output_path = f"{input_base}_slides.pptx"
    
    # 初始化PPT生成器
    ppt_generator = PPTGenerator(config)
    
    # 生成PPT
    output_file = ppt_generator.generate_ppt(content_analysis, output_path)
    
    print(f"PPT已生成: {output_file}")

if __name__ == "__main__":
    main()